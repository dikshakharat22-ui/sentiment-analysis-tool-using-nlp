from flask import Flask, render_template, request, jsonify
from textblob import TextBlob
import os
import logging
from werkzeug.utils import secure_filename
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from PIL import Image
import pytesseract
import io
import shutil

app = Flask(__name__)
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024
app.config['UPLOAD_FOLDER'] = '/tmp/uploads'

os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

Image.MAX_IMAGE_PIXELS = 178956970

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

ALLOWED_EXTENSIONS = {'txt', 'pdf', 'docx', 'pptx', 'png', 'jpg', 'jpeg'}

def check_tesseract_available():
    return shutil.which('tesseract') is not None

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_text_from_file(file, filename):
    ext = filename.rsplit('.', 1)[1].lower()
    
    try:
        if ext == 'txt':
            try:
                return file.read().decode('utf-8')
            except UnicodeDecodeError:
                file.seek(0)
                return file.read().decode('utf-8', errors='ignore')
        
        elif ext == 'pdf':
            pdf_reader = PdfReader(io.BytesIO(file.read()))
            text = ''
            for page in pdf_reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + '\n'
            return text.strip() if text else None
        
        elif ext == 'docx':
            doc = Document(io.BytesIO(file.read()))
            text = '\n'.join([paragraph.text for paragraph in doc.paragraphs])
            return text.strip() if text else None
        
        elif ext == 'pptx':
            prs = Presentation(io.BytesIO(file.read()))
            text = ''
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        text += shape.text + '\n'
            return text.strip() if text else None
        
        elif ext in ['png', 'jpg', 'jpeg']:
            if not check_tesseract_available():
                raise ValueError("OCR is not available. Please ensure Tesseract is installed on the system.")
            
            try:
                image = Image.open(io.BytesIO(file.read()))
                text = pytesseract.image_to_string(image)
                return text.strip() if text else None
            except Image.DecompressionBombError:
                raise ValueError("Image file is too large to process safely.")
        
        else:
            return None
            
    except ValueError as ve:
        raise ve
    except Exception as e:
        logger.error(f"Error extracting text from {ext} file: {str(e)}")
        raise ValueError(f"Unable to process {ext.upper()} file. Please ensure the file is not corrupted.")

def analyze_sentiment(text):
    if not text or not text.strip():
        return None
    
    blob = TextBlob(text)
    polarity = blob.sentiment.polarity
    subjectivity = blob.sentiment.subjectivity
    
    if polarity > 0.1:
        sentiment = 'Positive'
        emoji = '😊'
    elif polarity < -0.1:
        sentiment = 'Negative'
        emoji = '😞'
    else:
        sentiment = 'Neutral'
        emoji = '😐'
    
    return {
        'sentiment': sentiment,
        'polarity': round(polarity, 3),
        'subjectivity': round(subjectivity, 3),
        'emoji': emoji
    }

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/analyze', methods=['POST'])
def analyze():
    data = request.get_json()
    text = data.get('text', '')
    
    if not text:
        return jsonify({'error': 'No text provided'}), 400
    
    result = analyze_sentiment(text)
    if not result:
        return jsonify({'error': 'No valid text to analyze'}), 400
    
    return jsonify(result)

@app.route('/analyze-file', methods=['POST'])
def analyze_file():
    if 'file' not in request.files:
        return jsonify({'error': 'No file provided'}), 400
    
    file = request.files['file']
    
    if file.filename == '':
        return jsonify({'error': 'No file selected'}), 400
    
    if not file.filename or not allowed_file(file.filename):
        return jsonify({'error': f'File type not supported. Allowed types: TXT, PDF, DOCX, PPTX, PNG, JPG, JPEG'}), 400
    
    try:
        filename = secure_filename(file.filename) or 'uploaded_file'
        text = extract_text_from_file(file, filename)
        
        if not text or not text.strip():
            return jsonify({'error': 'No text could be extracted from the file. The file may be empty or contain only images.'}), 400
        
        result = analyze_sentiment(text)
        if not result:
            return jsonify({'error': 'Could not analyze the extracted text'}), 400
        
        result['extracted_text'] = text[:500] + '...' if len(text) > 500 else text
        result['filename'] = filename
        
        return jsonify(result)
    
    except ValueError as ve:
        logger.warning(f"User error: {str(ve)}")
        return jsonify({'error': str(ve)}), 400
    except Exception as e:
        logger.error(f"Unexpected error processing file: {str(e)}")
        return jsonify({'error': 'An error occurred while processing the file. Please try again with a different file.'}), 500

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=5000, debug=False)
